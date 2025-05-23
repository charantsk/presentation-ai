from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor as PPTXRGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from docx import Document
from fpdf import FPDF
import yaml
import io
from PIL import Image
import os
import logging
import re
import requests
import base64
from models import db, FileRecord
from services.image_service import (
    search_pexels_image,
    download_image,
    fetch_consistent_background_image
)
from config import Config
from jinja2 import Template

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# Pexels API configuration
PEXELS_API_KEY = Config.PEXELS_API_KEY
PEXELS_API_URL = "https://api.pexels.com/v1/search"

# HTML templates for different presentation styles
HTML_TEMPLATES = {
    'minimalist': """<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{{ title }}</title>
    <style>
        /* Reset and Base Styles */
        * {
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }
        body {
            font-family: 'Helvetica Neue', Arial, sans-serif;
            background-color: #f0f0f0;
            overflow-x: hidden;
        }

        /* Presentation Container */
        .presentation {
            width: 100vw;
            height: 100vh;
            background-color: #ffffff;
            position: relative;
        }

        /* Header Styles */
        .header {
            background-color: #ffffff;
            color: #333333;
            padding: 1.5rem;
            text-align: center;
            border-bottom: 1px solid #e0e0e0;
            position: fixed;
            top: 0;
            left: 0;
            right: 0;
            z-index: 10;
        }
        .header h1 {
            font-size: 2rem;
            font-weight: 300;
        }

        /* Slides Container */
        .slides-container {
            position: absolute;
            top: 5rem;
            width: 100%;
            height: calc(100vh - 5rem);
            overflow: hidden;
        }
        .slides {
            display: flex;
            height: 100%;
            transition: transform 0.5s ease;
        }

        /* Slide Styles */
        .slide {
            min-width: 100%;
            height: 100%;
            display: flex;
            flex-direction: row;
            align-items: center;
            padding: 2rem;
        }
        .slide-image-container {
            width: 40%;
            padding: 1rem;
        }
        .slide-image {
            max-width: 100%;
            max-height: 30vh;
            object-fit: cover;
            border-radius: 0.5rem;
        }
        .slide-content-container {
            width: 60%;
            padding: 1rem;
        }
        .slide-title {
            font-size: 1.8rem;
            font-weight: 500;
            color: #333333;
            margin-bottom: 1rem;
        }
        .slide-content {
            font-size: 1.2rem;
            color: #555555;
            margin-bottom: 1rem;
        }
        .bullets {
            list-style-type: disc;
            padding-left: 1.5rem;
        }
        .bullets li {
            font-size: 1.1rem;
            color: #555555;
            margin-bottom: 0.5rem;
        }

        /* Navigation Styles */
        .navigation {
            position: fixed;
            bottom: 1.5rem;
            left: 50%;
            transform: translateX(-50%);
            display: flex;
            gap: 1rem;
            z-index: 10;
        }
        .nav-btn {
            background-color: #333333;
            color: #ffffff;
            border: none;
            border-radius: 0.3rem;
            width: 2.5rem;
            height: 2.5rem;
            font-size: 1.2rem;
            cursor: pointer;
            display: flex;
            align-items: center;
            justify-content: center;
        }
        .nav-btn:hover {
            background-color: #555555;
        }
        .nav-btn:disabled {
            background-color: #cccccc;
            cursor: not-allowed;
        }

        /* Slide Indicator and Progress Bar */
        .slide-indicator {
            position: fixed;
            bottom: 1.5rem;
            right: 1.5rem;
            background-color: #333333;
            color: #ffffff;
            padding: 0.4rem 0.8rem;
            border-radius: 1rem;
            font-size: 0.9rem;
        }
        .progress-bar {
            position: fixed;
            bottom: 0;
            left: 0;
            height: 0.25rem;
            background-color: #333333;
            transition: width 0.3s;
        }

        /* Responsive Design */
        @media screen and (max-width: 768px) {
            .slide {
                flex-direction: column;
            }
            .slide-image-container,
            .slide-content-container {
                width: 100%;
            }
            .slide-image {
                max-height: 20vh;
            }
            .slide-title {
                font-size: 1.5rem;
            }
            .slide-content,
            .bullets li {
                font-size: 1rem;
            }
            .header h1 {
                font-size: 1.5rem;
            }
        }
    </style>
</head>
<body>
    <div class="presentation">
        <header class="header">
            <h1>{{ title }}</h1>
        </header>
        <div class="slides-container">
            <div class="slides" id="slides">
                {% for slide in slides %}
                <section class="slide">
                    {% if slide.image %}
                    <div class="slide-image-container">
                        <img class="slide-image" src="{{ slide.image }}" alt="{{ slide.title }}">
                    </div>
                    {% endif %}
                    <div class="slide-content-container">
                        <h2 class="slide-title">{{ slide.title }}</h2>
                        {% if slide.content %}
                        <p class="slide-content">{{ slide.content }}</p>
                        {% endif %}
                        {% if slide.bullets %}
                        <ul class="bullets">
                            {% for bullet in slide.bullets %}
                            <li>{{ bullet }}</li>
                            {% endfor %}
                        </ul>
                        {% endif %}
                    </div>
                </section>
                {% endfor %}
            </div>
        </div>
        <nav class="navigation">
            <button class="nav-btn" id="prevBtn" aria-label="Previous Slide">←</button>
            <button class="nav-btn" id="nextBtn" aria-label="Next Slide">→</button>
        </nav>
        <div class="slide-indicator" id="slideIndicator" aria-live="polite">1 / {{ slides|length }}</div>
        <div class="progress-bar" id="progressBar" role="progressbar"></div>
    </div>

    <script>
        document.addEventListener('DOMContentLoaded', () => {
            const slides = document.getElementById('slides');
            const prevBtn = document.getElementById('prevBtn');
            const nextBtn = document.getElementById('nextBtn');
            const slideIndicator = document.getElementById('slideIndicator');
            const progressBar = document.getElementById('progressBar');
            const totalSlides = {{ slides|length }};
            let currentSlide = 0;

            const updateSlide = () => {
                slides.style.transform = `translateX(-${currentSlide * 100}%)`;
                slideIndicator.textContent = `${currentSlide + 1} / ${totalSlides}`;
                prevBtn.disabled = currentSlide === 0;
                nextBtn.disabled = currentSlide === totalSlides - 1;
                progressBar.style.width = `${((currentSlide + 1) / totalSlides) * 100}%`;
            };

            const previousSlide = () => {
                if (currentSlide > 0) {
                    currentSlide--;
                    updateSlide();
                }
            };

            const nextSlide = () => {
                if (currentSlide < totalSlides - 1) {
                    currentSlide++;
                    updateSlide();
                }
            };

            prevBtn.addEventListener('click', previousSlide);
            nextBtn.addEventListener('click', nextSlide);
            document.addEventListener('keydown', (e) => {
                if (e.key === 'ArrowLeft') previousSlide();
                if (e.key === 'ArrowRight') nextSlide();
            });
            window.addEventListener('resize', updateSlide);

            updateSlide();
        });
    </script>
</body>
</html>
""",
    'modern': """<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{{ title }}</title>
    <style>
        /* Reset and Base Styles */
        * {
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }
        body {
            font-family: 'Roboto', sans-serif;
            background-color: #e3f2fd;
            overflow-x: hidden;
        }

        /* Presentation Container */
        .presentation {
            width: 100vw;
            height: 100vh;
            background-color: #ffffff;
            position: relative;
            box-shadow: 0 0.25rem 0.75rem rgba(0, 0, 0, 0.1);
        }

        /* Header Styles */
        .header {
            background: linear-gradient(90deg, #0288d1, #4fc3f7);
            color: #ffffff;
            padding: 1rem;
            text-align: center;
            position: fixed;
            top: 0;
            left: 0;
            right: 0;
            z-index: 10;
        }
        .header h1 {
            font-size: 2.2rem;
            font-weight: 500;
        }

        /* Slides Container */
        .slides-container {
            position: absolute;
            top: 4.5rem;
            width: 100%;
            height: calc(100vh - 4.5rem);
            overflow: hidden;
        }
        .slides {
            display: flex;
            height: 100%;
            transition: transform 0.5s ease;
        }

        /* Slide Styles */
        .slide {
            min-width: 100%;
            height: 100%;
            display: flex;
            flex-direction: row;
            align-items: center;
            padding: 2rem;
            background-color: #fafafa;
        }
        .slide-image-container {
            width: 40%;
            padding: 1rem;
        }
        .slide-image {
            max-width: 100%;
            max-height: 30vh;
            object-fit: cover;
            border-radius: 0.625rem;
            box-shadow: 0 0.125rem 0.5rem rgba(0, 0, 0, 0.15);
        }
        .slide-content-container {
            width: 60%;
            padding: 1rem;
        }
        .slide-title {
            font-size: 2rem;
            font-weight: 600;
            color: #0277bd;
            margin-bottom: 1.25rem;
        }
        .slide-content {
            font-size: 1.3rem;
            color: #424242;
            margin-bottom: 1.25rem;
        }
        .bullets {
            list-style-type: square;
            padding-left: 1.5rem;
        }
        .bullets li {
            font-size: 1.2rem;
            color: #424242;
            margin-bottom: 0.75rem;
        }

        /* Navigation Styles */
        .navigation {
            position: fixed;
            bottom: 1.5rem;
            left: 50%;
            transform: translateX(-50%);
            display: flex;
            gap: 1.25rem;
            z-index: 10;
        }
        .nav-btn {
            background-color: #0288d1;
            color: #ffffff;
            border: none;
            border-radius: 50%;
            width: 3rem;
            height: 3rem;
            font-size: 1.25rem;
            cursor: pointer;
            display: flex;
            align-items: center;
            justify-content: center;
        }
        .nav-btn:hover {
            background-color: #0277bd;
        }
        .nav-btn:disabled {
            background-color: #b0bec5;
            cursor: not-allowed;
        }

        /* Slide Indicator and Progress Bar */
        .slide-indicator {
            position: fixed;
            bottom: 1.5rem;
            right: 1.5rem;
            background-color: #0288d1;
            color: #ffffff;
            padding: 0.5rem 0.75rem;
            border-radius: 1.25rem;
            font-size: 0.875rem;
        }
        .progress-bar {
            position: fixed;
            bottom: 0;
            left: 0;
            height: 0.3125rem;
            background-color: #4fc3f7;
            transition: width 0.3s;
        }

        /* Responsive Design */
        @media screen and (max-width: 768px) {
            .slide {
                flex-direction: column;
            }
            .slide-image-container,
            .slide-content-container {
                width: 100%;
            }
            .slide-image {
                max-height: 20vh;
            }
            .slide-title {
                font-size: 1.6rem;
            }
            .slide-content,
            .bullets li {
                font-size: 1.1rem;
            }
            .header h1 {
                font-size: 1.8rem;
            }
        }
    </style>
</head>
<body>
    <div class="presentation">
        <header class="header">
            <h1>{{ title }}</h1>
        </header>
        <div class="slides-container">
            <div class="slides" id="slides">
                {% for slide in slides %}
                <section class="slide">
                    {% if slide.image %}
                    <div class="slide-image-container">
                        <img class="slide-image" src="{{ slide.image }}" alt="{{ slide.title }}">
                    </div>
                    {% endif %}
                    <div class="slide-content-container">
                        <h2 class="slide-title">{{ slide.title }}</h2>
                        {% if slide.content %}
                        <p class="slide-content">{{ slide.content }}</p>
                        {% endif %}
                        {% if slide.bullets %}
                        <ul class="bullets">
                            {% for bullet in slide.bullets %}
                            <li>{{ bullet }}</li>
                            {% endfor %}
                        </ul>
                        {% endif %}
                    </div>
                </section>
                {% endfor %}
            </div>
        </div>
        <nav class="navigation">
            <button class="nav-btn" id="prevBtn" aria-label="Previous Slide">←</button>
            <button class="nav-btn" id="nextBtn" aria-label="Next Slide">→</button>
        </nav>
        <div class="slide-indicator" id="slideIndicator" aria-live="polite">1 / {{ slides|length }}</div>
        <div class="progress-bar" id="progressBar" role="progressbar"></div>
    </div>

    <script>
        document.addEventListener('DOMContentLoaded', () => {
            const slides = document.getElementById('slides');
            const prevBtn = document.getElementById('prevBtn');
            const nextBtn = document.getElementById('nextBtn');
            const slideIndicator = document.getElementById('slideIndicator');
            const progressBar = document.getElementById('progressBar');
            const totalSlides = {{ slides|length }};
            let currentSlide = 0;

            const updateSlide = () => {
                slides.style.transform = `translateX(-${currentSlide * 100}%)`;
                slideIndicator.textContent = `${currentSlide + 1} / ${totalSlides}`;
                prevBtn.disabled = currentSlide === 0;
                nextBtn.disabled = currentSlide === totalSlides - 1;
                progressBar.style.width = `${((currentSlide + 1) / totalSlides) * 100}%`;
            };

            const previousSlide = () => {
                if (currentSlide > 0) {
                    currentSlide--;
                    updateSlide();
                }
            };

            const nextSlide = () => {
                if (currentSlide < totalSlides - 1) {
                    currentSlide++;
                    updateSlide();
                }
            };

            prevBtn.addEventListener('click', previousSlide);
            nextBtn.addEventListener('click', nextSlide);
            document.addEventListener('keydown', (e) => {
                if (e.key === 'ArrowLeft') previousSlide();
                if (e.key === 'ArrowRight') nextSlide();
            });
            window.addEventListener('resize', updateSlide);

            updateSlide();
        });
    </script>
</body>
</html>
""",
    'professional': """<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{{ title }}</title>
    <style>
        /* Reset and Base Styles */
        * {
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }
        body {
            font-family: 'Arial', sans-serif;
            background-color: #eceff1;
            overflow-x: hidden;
        }

        /* Presentation Container */
        .presentation {
            width: 100vw;
            height: 100vh;
            background-color: #ffffff;
            position: relative;
        }

        /* Header Styles */
        .header {
            background-color: #263238;
            color: #ffffff;
            padding: 1.5rem;
            text-align: center;
            position: fixed;
            top: 0;
            left: 0;
            right: 0;
            z-index: 10;
        }
        .header h1 {
            font-size: 2rem;
            font-weight: 400;
        }

        /* Slides Container */
        .slides-container {
            position: absolute;
            top: 5rem;
            width: 100%;
            height: calc(100vh - 5rem);
            overflow: hidden;
        }
        .slides {
            display: flex;
            height: 100%;
            transition: transform 0.5s ease;
        }

        /* Slide Styles */
        .slide {
            min-width: 100%;
            height: 100%;
            display: flex;
            flex-direction: row;
            align-items: center;
            padding: 2rem;
            background-color: #ffffff;
        }
        .slide-image-container {
            width: 40%;
            padding: 1rem;
        }
        .slide-image {
            max-width: 100%;
            max-height: 30vh;
            object-fit: cover;
            border: 1px solid #e0e0e0;
        }
        .slide-content-container {
            width: 60%;
            padding: 1rem;
        }
        .slide-title {
            font-size: 1.9rem;
            font-weight: 500;
            color: #263238;
            margin-bottom: 1rem;
        }
        .slide-content {
            font-size: 1.2rem;
            color: #37474f;
            margin-bottom: 1rem;
        }
        .bullets {
            list-style-type: circle;
            padding-left: 1.5rem;
        }
        .bullets li {
            font-size: 1.1rem;
            color: #37474f;
            margin-bottom: 0.5rem;
        }

        /* Navigation Styles */
        .navigation {
            position: fixed;
            bottom: 1.5rem;
            left: 50%;
            transform: translateX(-50%);
            display: flex;
            gap: 1rem;
            z-index: 10;
        }
        .nav-btn {
            background-color: #263238;
            color: #ffffff;
            border: none;
            border-radius: 0.3rem;
            width: 2.5rem;
            height: 2.5rem;
            font-size: 1.2rem;
            cursor: pointer;
            display: flex;
            align-items: center;
            justify-content: center;
        }
        .nav-btn:hover {
            background-color: #37474f;
        }
        .nav-btn:disabled {
            background-color: #b0bec5;
            cursor: not-allowed;
        }

        /* Slide Indicator and Progress Bar */
        .slide-indicator {
            position: fixed;
            bottom: 1.5rem;
            right: 1.5rem;
            background-color: #263238;
            color: #ffffff;
            padding: 0.4rem 0.8rem;
            border-radius: 1rem;
            font-size: 0.9rem;
        }
        .progress-bar {
            position: fixed;
            bottom: 0;
            left: 0;
            height: 0.25rem;
            background-color: #263238;
            transition: width 0.3s;
        }

        /* Responsive Design */
        @media screen and (max-width: 768px) {
            .slide {
                flex-direction: column;
            }
            .slide-image-container,
            .slide-content-container {
                width: 100%;
            }
            .slide-image {
                max-height: 20vh;
            }
            .slide-title {
                font-size: 1.5rem;
            }
            .slide-content,
            .bullets li {
                font-size: 1rem;
            }
            .header h1 {
                font-size: 1.5rem;
            }
        }
    </style>
</head>
<body>
    <div class="presentation">
        <header class="header">
            <h1>{{ title }}</h1>
        </header>
        <div class="slides-container">
            <div class="slides" id="slides">
                {% for slide in slides %}
                <section class="slide">
                    {% if slide.image %}
                    <div class="slide-image-container">
                        <img class="slide-image" src="{{ slide.image }}" alt="{{ slide.title }}">
                    </div>
                    {% endif %}
                    <div class="slide-content-container">
                        <h2 class="slide-title">{{ slide.title }}</h2>
                        {% if slide.content %}
                        <p class="slide-content">{{ slide.content }}</p>
                        {% endif %}
                        {% if slide.bullets %}
                        <ul class="bullets">
                            {% for bullet in slide.bullets %}
                            <li>{{ bullet }}</li>
                            {% endfor %}
                        </ul>
                        {% endif %}
                    </div>
                </section>
                {% endfor %}
            </div>
        </div>
        <nav class="navigation">
            <button class="nav-btn" id="prevBtn" aria-label="Previous Slide">←</button>
            <button class="nav-btn" id="nextBtn" aria-label="Next Slide">→</button>
        </nav>
        <div class="slide-indicator" id="slideIndicator" aria-live="polite">1 / {{ slides|length }}</div>
        <div class="progress-bar" id="progressBar" role="progressbar"></div>
    </div>

    <script>
        document.addEventListener('DOMContentLoaded', () => {
            const slides = document.getElementById('slides');
            const prevBtn = document.getElementById('prevBtn');
            const nextBtn = document.getElementById('nextBtn');
            const slideIndicator = document.getElementById('slideIndicator');
            const progressBar = document.getElementById('progressBar');
            const totalSlides = {{ slides|length }};
            let currentSlide = 0;

            const updateSlide = () => {
                slides.style.transform = `translateX(-${currentSlide * 100}%)`;
                slideIndicator.textContent = `${currentSlide + 1} / ${totalSlides}`;
                prevBtn.disabled = currentSlide === 0;
                nextBtn.disabled = currentSlide === totalSlides - 1;
                progressBar.style.width = `${((currentSlide + 1) / totalSlides) * 100}%`;
            };

            const previousSlide = () => {
                if (currentSlide > 0) {
                    currentSlide--;
                    updateSlide();
                }
            };

            const nextSlide = () => {
                if (currentSlide < totalSlides - 1) {
                    currentSlide++;
                    updateSlide();
                }
            };

            prevBtn.addEventListener('click', previousSlide);
            nextBtn.addEventListener('click', nextSlide);
            document.addEventListener('keydown', (e) => {
                if (e.key === 'ArrowLeft') previousSlide();
                if (e.key === 'ArrowRight') nextSlide();
            });
            window.addEventListener('resize', updateSlide);

            updateSlide();
        });
    </script>
</body>
</html>
""",
    'corporate': """<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{{ title }}</title>
    <style>
        /* Reset and Base Styles */
        * {
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }
        body {
            font-family: 'Inter', 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            background: linear-gradient(135deg, #f8fafc 0%, #e2e8f0 100%);
            overflow-x: hidden;
            line-height: 1.6;
        }

        /* Presentation Container */
        .presentation {
            width: 100vw;
            height: 100vh;
            background-color: #ffffff;
            position: relative;
            box-shadow: 0 20px 40px rgba(0, 0, 0, 0.1);
        }

        /* Header Styles */
        .header {
            background: linear-gradient(135deg, #1e293b 0%, #334155 100%);
            color: #ffffff;
            padding: 2rem 3rem;
            text-align: center;
            position: fixed;
            top: 0;
            left: 0;
            right: 0;
            z-index: 10;
            border-bottom: 4px solid #0ea5e9;
            box-shadow: 0 4px 20px rgba(0, 0, 0, 0.15);
        }
        .header h1 {
            font-size: 2.5rem;
            font-weight: 600;
            letter-spacing: -0.025em;
            color: #ffffff;
        }
        .header::after {
            content: '';
            position: absolute;
            bottom: -4px;
            left: 50%;
            transform: translateX(-50%);
            width: 100px;
            height: 4px;
            background: linear-gradient(90deg, #0ea5e9, #06b6d4);
        }

        /* Slides Container */
        .slides-container {
            position: absolute;
            top: 6rem;
            width: 100%;
            height: calc(100vh - 6rem);
            overflow: hidden;
        }
        .slides {
            display: flex;
            height: 100%;
            transition: transform 0.6s cubic-bezier(0.4, 0, 0.2, 1);
        }

        /* Slide Styles */
        .slide {
            min-width: 100%;
            height: 100%;
            display: flex;
            flex-direction: row;
            align-items: center;
            padding: 3rem 4rem;
            background: linear-gradient(135deg, #ffffff 0%, #f8fafc 100%);
            position: relative;
        }
        .slide::before {
            content: '';
            position: absolute;
            top: 0;
            left: 0;
            right: 0;
            height: 2px;
            background: linear-gradient(90deg, #0ea5e9, #06b6d4, #8b5cf6);
        }
        .slide-image-container {
            width: 45%;
            padding: 2rem;
            display: flex;
            justify-content: center;
            align-items: center;
        }
        .slide-image {
            max-width: 100%;
            max-height: 50vh;
            object-fit: contain;
            border-radius: 12px;
            box-shadow: 0 10px 30px rgba(0, 0, 0, 0.1);
            border: 1px solid #e2e8f0;
        }
        .slide-content-container {
            width: 55%;
            padding: 2rem 3rem;
        }
        .slide-title {
            font-size: 2.25rem;
            font-weight: 700;
            color: #1e293b;
            margin-bottom: 1.5rem;
            line-height: 1.2;
            position: relative;
        }
        .slide-title::after {
            content: '';
            position: absolute;
            bottom: -8px;
            left: 0;
            width: 60px;
            height: 3px;
            background: linear-gradient(90deg, #0ea5e9, #06b6d4);
            border-radius: 2px;
        }
        .slide-content {
            font-size: 1.2rem;
            color: #475569;
            margin-bottom: 2rem;
            line-height: 1.7;
        }
        .bullets {
            list-style: none;
            padding-left: 0;
        }
        .bullets li {
            font-size: 1.1rem;
            color: #475569;
            margin-bottom: 1.2rem;
            padding: 1rem 1.5rem;
            background: #f8fafc;
            border-radius: 8px;
            border-left: 4px solid #0ea5e9;
            box-shadow: 0 2px 8px rgba(0, 0, 0, 0.05);
            transition: all 0.3s ease;
            position: relative;
        }
        .bullets li::before {
            content: '▶';
            position: absolute;
            left: 0.75rem;
            top: 50%;
            transform: translateY(-50%);
            color: #0ea5e9;
            font-size: 0.8rem;
        }
        .bullets li:hover {
            transform: translateX(8px);
            box-shadow: 0 4px 12px rgba(0, 0, 0, 0.1);
            background: #ffffff;
        }

        /* Navigation Styles */
        .navigation {
            position: fixed;
            bottom: 2rem;
            left: 50%;
            transform: translateX(-50%);
            display: flex;
            gap: 1rem;
            z-index: 10;
        }
        .nav-btn {
            background: linear-gradient(135deg, #1e293b 0%, #334155 100%);
            color: #ffffff;
            border: none;
            border-radius: 8px;
            width: 3rem;
            height: 3rem;
            font-size: 1.2rem;
            cursor: pointer;
            display: flex;
            align-items: center;
            justify-content: center;
            box-shadow: 0 4px 12px rgba(30, 41, 59, 0.3);
            transition: all 0.3s ease;
        }
        .nav-btn:hover {
            background: linear-gradient(135deg, #334155 0%, #475569 100%);
            transform: translateY(-2px);
            box-shadow: 0 6px 16px rgba(30, 41, 59, 0.4);
        }
        .nav-btn:disabled {
            background: #cbd5e1;
            cursor: not-allowed;
            transform: none;
            box-shadow: none;
        }

        /* Slide Indicator and Progress Bar */
        .slide-indicator {
            position: fixed;
            bottom: 2rem;
            right: 2rem;
            background: linear-gradient(135deg, #1e293b 0%, #334155 100%);
            color: #ffffff;
            padding: 0.75rem 1.25rem;
            border-radius: 20px;
            font-size: 0.9rem;
            font-weight: 600;
            box-shadow: 0 4px 12px rgba(30, 41, 59, 0.3);
        }
        .progress-bar {
            position: fixed;
            bottom: 0;
            left: 0;
            height: 4px;
            background: linear-gradient(90deg, #0ea5e9, #06b6d4, #8b5cf6);
            transition: width 0.4s ease;
            box-shadow: 0 -2px 8px rgba(14, 165, 233, 0.3);
        }

        /* Responsive Design */
        @media screen and (max-width: 1024px) {
            .slide {
                padding: 2rem 3rem;
            }
            .slide-title {
                font-size: 2rem;
            }
            .slide-content {
                font-size: 1.1rem;
            }
        }

        @media screen and (max-width: 768px) {
            .header {
                padding: 1.5rem;
            }
            .header h1 {
                font-size: 1.8rem;
            }
            .slides-container {
                top: 5rem;
                height: calc(100vh - 5rem);
            }
            .slide {
                flex-direction: column;
                padding: 2rem 1.5rem;
                text-align: center;
            }
            .slide-image-container,
            .slide-content-container {
                width: 100%;
                padding: 1rem;
            }
            .slide-image {
                max-height: 30vh;
                margin-bottom: 1rem;
            }
            .slide-title {
                font-size: 1.75rem;
            }
            .slide-content,
            .bullets li {
                font-size: 1rem;
            }
            .bullets li {
                padding: 0.8rem 1rem;
                margin-bottom: 1rem;
            }
            .navigation {
                bottom: 1rem;
            }
            .nav-btn {
                width: 2.5rem;
                height: 2.5rem;
                font-size: 1rem;
            }
            .slide-indicator {
                bottom: 1rem;
                right: 1rem;
                padding: 0.5rem 1rem;
                font-size: 0.8rem;
            }
        }

        @media screen and (max-width: 480px) {
            .header {
                padding: 1rem;
            }
            .header h1 {
                font-size: 1.5rem;
            }
            .slide {
                padding: 1.5rem 1rem;
            }
            .slide-title {
                font-size: 1.5rem;
            }
            .slide-content {
                font-size: 0.95rem;
            }
            .bullets li {
                font-size: 0.9rem;
                padding: 0.7rem 0.8rem;
            }
        }
    </style>
</head>
<body>
    <div class="presentation">
        <header class="header">
            <h1>{{ title }}</h1>
        </header>
        <div class="slides-container">
            <div class="slides" id="slides">
                {% for slide in slides %}
                <section class="slide">
                    {% if slide.image %}
                    <div class="slide-image-container">
                        <img class="slide-image" src="{{ slide.image }}" alt="{{ slide.title }}">
                    </div>
                    {% endif %}
                    <div class="slide-content-container">
                        <h2 class="slide-title">{{ slide.title }}</h2>
                        {% if slide.content %}
                        <p class="slide-content">{{ slide.content }}</p>
                        {% endif %}
                        {% if slide.bullets %}
                        <ul class="bullets">
                            {% for bullet in slide.bullets %}
                            <li>{{ bullet }}</li>
                            {% endfor %}
                        </ul>
                        {% endif %}
                    </div>
                </section>
                {% endfor %}
            </div>
        </div>
        <nav class="navigation">
            <button class="nav-btn" id="prevBtn" aria-label="Previous Slide">←</button>
            <button class="nav-btn" id="nextBtn" aria-label="Next Slide">→</button>
        </nav>
        <div class="slide-indicator" id="slideIndicator" aria-live="polite">1 / {{ slides|length }}</div>
        <div class="progress-bar" id="progressBar" role="progressbar"></div>
    </div>

    <script>
        document.addEventListener('DOMContentLoaded', () => {
            const slides = document.getElementById('slides');
            const prevBtn = document.getElementById('prevBtn');
            const nextBtn = document.getElementById('nextBtn');
            const slideIndicator = document.getElementById('slideIndicator');
            const progressBar = document.getElementById('progressBar');
            const totalSlides = {{ slides|length }};
            let currentSlide = 0;

            const updateSlide = () => {
                slides.style.transform = `translateX(-${currentSlide * 100}%)`;
                slideIndicator.textContent = `${currentSlide + 1} / ${totalSlides}`;
                prevBtn.disabled = currentSlide === 0;
                nextBtn.disabled = currentSlide === totalSlides - 1;
                progressBar.style.width = `${((currentSlide + 1) / totalSlides) * 100}%`;
            };

            const previousSlide = () => {
                if (currentSlide > 0) {
                    currentSlide--;
                    updateSlide();
                }
            };

            const nextSlide = () => {
                if (currentSlide < totalSlides - 1) {
                    currentSlide++;
                    updateSlide();
                }
            };

            prevBtn.addEventListener('click', previousSlide);
            nextBtn.addEventListener('click', nextSlide);
            document.addEventListener('keydown', (e) => {
                if (e.key === 'ArrowLeft') previousSlide();
                if (e.key === 'ArrowRight') nextSlide();
            });
            window.addEventListener('resize', updateSlide);

            updateSlide();
        });
    </script>
</body>
</html>
""",
    'executive': """<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{{ title }}</title>
    <style>
        /* Reset and Base Styles */
        * {
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }
        body {
            font-family: 'Playfair Display', 'Georgia', serif;
            background: linear-gradient(135deg, #0f172a 0%, #1e293b 50%, #334155 100%);
            overflow-x: hidden;
            line-height: 1.6;
        }

        /* Presentation Container */
        .presentation {
            width: 100vw;
            height: 100vh;
            background: linear-gradient(135deg, #ffffff 0%, #f8fafc 100%);
            position: relative;
            border: 2px solid #e2e8f0;
        }

        /* Header Styles */
        .header {
            background: linear-gradient(135deg, #0f172a 0%, #1e293b 100%);
            color: #f8fafc;
            padding: 3rem;
            text-align: center;
            position: fixed;
            top: 0;
            left: 0;
            right: 0;
            z-index: 10;
            border-bottom: 6px solid #d4af37;
            box-shadow: 0 8px 32px rgba(15, 23, 42, 0.4);
        }
        .header h1 {
            font-size: 3rem;
            font-weight: 700;
            letter-spacing: 2px;
            color: #f8fafc;
            position: relative;
            text-shadow: 2px 2px 4px rgba(0, 0, 0, 0.3);
        }
        .header::after {
            content: '';
            position: absolute;
            bottom: -6px;
            left: 50%;
            transform: translateX(-50%);
            width: 120px;
            height: 6px;
            background: linear-gradient(90deg, #d4af37, #f7931e, #d4af37);
        }

        /* Slides Container */
        .slides-container {
            position: absolute;
            top: 8rem;
            width: 100%;
            height: calc(100vh - 8rem);
            overflow: hidden;
        }
        .slides {
            display: flex;
            height: 100%;
            transition: transform 0.8s cubic-bezier(0.25, 0.46, 0.45, 0.94);
        }

        /* Slide Styles */
        .slide {
            min-width: 100%;
            height: 100%;
            display: flex;
            flex-direction: row;
            align-items: center;
            padding: 4rem 5rem;
            background: linear-gradient(135deg, #ffffff 0%, #f8fafc 100%);
            position: relative;
        }
        .slide::before {
            content: '';
            position: absolute;
            top: 0;
            left: 0;
            right: 0;
            bottom: 0;
            background: 
                linear-gradient(90deg, transparent 0%, rgba(212, 175, 55, 0.05) 50%, transparent 100%),
                linear-gradient(0deg, transparent 0%, rgba(212, 175, 55, 0.03) 50%, transparent 100%);
            pointer-events: none;
        }
        .slide::after {
            content: '';
            position: absolute;
            top: 2rem;
            left: 2rem;
            right: 2rem;
            bottom: 2rem;
            border: 2px solid rgba(212, 175, 55, 0.2);
            border-radius: 8px;
            pointer-events: none;
        }
        .slide-image-container {
            width: 40%;
            padding: 2rem;
            z-index: 2;
            display: flex;
            justify-content: center;
            align-items: center;
        }
        .slide-image {
            max-width: 100%;
            max-height: 50vh;
            object-fit: contain;
            border-radius: 8px;
            box-shadow: 0 12px 32px rgba(0, 0, 0, 0.15);
            border: 2px solid #e2e8f0;
        }
        .slide-content-container {
            width: 60%;
            padding: 2rem 3rem;
            z-index: 2;
        }
        .slide-title {
            font-size: 2.75rem;
            font-weight: 700;
            color: #0f172a;
            margin-bottom: 2rem;
            letter-spacing: 1px;
            line-height: 1.2;
            position: relative;
            font-style: italic;
        }
        .slide-title::after {
            content: '';
            position: absolute;
            bottom: -10px;
            left: 0;
            width: 80px;
            height: 4px;
            background: linear-gradient(90deg, #d4af37, #f7931e);
            border-radius: 2px;
        }
        .slide-content {
            font-size: 1.3rem;
            color: #475569;
            margin-bottom: 2rem;
            line-height: 1.8;
            font-family: 'Inter', 'Segoe UI', sans-serif;
        }
        .bullets {
            list-style: none;
            padding-left: 0;
        }
        .bullets li {
            font-size: 1.2rem;
            color: #334155;
            margin-bottom: 1.5rem;
            padding: 1.2rem 2rem;
            background: linear-gradient(135deg, #f8fafc 0%, #ffffff 100%);
            border-radius: 12px;
            border-left: 5px solid #d4af37;
            box-shadow: 0 4px 16px rgba(0, 0, 0, 0.08);
            position: relative;
            transition: all 0.4s ease;
            font-family: 'Inter', 'Segoe UI', sans-serif;
            line-height: 1.6;
        }
        .bullets li::before {
            content: '◆';
            position: absolute;
            left: 1rem;
            top: 50%;
            transform: translateY(-50%);
            color: #d4af37;
            font-size: 1rem;
        }
        .bullets li:hover {
            transform: translateX(12px);
            box-shadow: 0 8px 24px rgba(0, 0, 0, 0.12);
            background: linear-gradient(135deg, #ffffff 0%, #f8fafc 100%);
            border-left-width: 6px;
        }

        /* Navigation Styles */
        .navigation {
            position: fixed;
            bottom: 2.5rem;
            left: 50%;
            transform: translateX(-50%);
            display: flex;
            gap: 1.5rem;
            z-index: 10;
        }
        .nav-btn {
            background: linear-gradient(135deg, #0f172a 0%, #1e293b 100%);
            color: #f8fafc;
            border: 2px solid #d4af37;
            border-radius: 12px;
            width: 3.5rem;
            height: 3.5rem;
            font-size: 1.4rem;
            cursor: pointer;
            display: flex;
            align-items: center;
            justify-content: center;
            box-shadow: 0 6px 20px rgba(15, 23, 42, 0.4);
            transition: all 0.4s ease;
            font-weight: 600;
        }
        .nav-btn:hover {
            background: linear-gradient(135deg, #1e293b 0%, #334155 100%);
            transform: translateY(-3px) scale(1.05);
            box-shadow: 0 8px 28px rgba(15, 23, 42, 0.5);
            border-color: #f7931e;
        }
        .nav-btn:disabled {
            background: #cbd5e1;
            color: #94a3b8;
            border-color: #cbd5e1;
            cursor: not-allowed;
            transform: none;
            box-shadow: 0 2px 8px rgba(0, 0, 0, 0.1);
        }

        /* Slide Indicator and Progress Bar */
        .slide-indicator {
            position: fixed;
            bottom: 2.5rem;
            right: 2.5rem;
            background: linear-gradient(135deg, #0f172a 0%, #1e293b 100%);
            color: #f8fafc;
            border: 2px solid #d4af37;
            padding: 1rem 1.5rem;
            border-radius: 20px;
            font-size: 1rem;
            font-weight: 700;
            box-shadow: 0 6px 20px rgba(15, 23, 42, 0.4);
            letter-spacing: 1px;
        }
        .progress-bar {
            position: fixed;
            bottom: 0;
            left: 0;
            height: 6px;
            background: linear-gradient(90deg, #d4af37 0%, #f7931e 50%, #d4af37 100%);
            transition: width 0.6s ease;
            box-shadow: 0 -3px 12px rgba(212, 175, 55, 0.4);
        }

        /* Responsive Design */
        @media screen and (max-width: 1024px) {
            .slide {
                padding: 3rem 4rem;
            }
            .slide-title {
                font-size: 2.25rem;
            }
            .slide-content {
                font-size: 1.2rem;
            }
            .bullets li {
                font-size: 1.1rem;
                padding: 1rem 1.8rem;
            }
        }

        @media screen and (max-width: 768px) {
            .header {
                padding: 2rem 1.5rem;
            }
            .header h1 {
                font-size: 2.2rem;
                letter-spacing: 1px;
            }
            .slides-container {
                top: 7rem;
                height: calc(100vh - 7rem);
            }
            .slide {
                flex-direction: column;
                padding: 2rem 1.5rem;
                text-align: center;
            }
            .slide-image-container,
            .slide-content-container {
                width: 100%;
                padding: 1rem;
            }
            .slide-image {
                max-height: 35vh;
                margin-bottom: 1.5rem;
            }
            .slide-title {
                font-size: 2rem;
            }
            .slide-content {
                font-size: 1.1rem;
            }
            .bullets li {
                font-size: 1rem;
                padding: 0.9rem 1.5rem;
                margin-bottom: 1.2rem;
            }
            .navigation {
                bottom: 1.5rem;
                gap: 1rem;
            }
            .nav-btn {
                width: 3rem;
                height: 3rem;
                font-size: 1.2rem;
            }
            .slide-indicator {
                bottom: 1.5rem;
                right: 1.5rem;
                padding: 0.8rem 1.2rem;
                font-size: 0.9rem;
            }
        }

        @media screen and (max-width: 480px) {
            .header {
                padding: 1.5rem 1rem;
            }
            .header h1 {
                font-size: 1.8rem;
                letter-spacing: 0.5px;
            }
            .slides-container {
                top: 6rem;
                height: calc(100vh - 6rem);
            }
            .slide {
                padding: 1.5rem 1rem;
            }
            .slide-title {
                font-size: 1.75rem;
            }
            .slide-content {
                font-size: 1rem;
            }
            .bullets li {
                font-size: 0.95rem;
                padding: 0.8rem 1.2rem;
            }
            .nav-btn {
                width: 2.8rem;
                height: 2.8rem;
                font-size: 1.1rem;
            }
            .slide-indicator {
                padding: 0.7rem 1rem;
                font-size: 0.8rem;
            }
        }
    </style>
</head>
<body>
    <div class="presentation">
        <header class="header">
            <h1>{{ title }}</h1>
        </header>
        <div class="slides-container">
            <div class="slides" id="slides">
                {% for slide in slides %}
                <section class="slide">
                    {% if slide.image %}
                    <div class="slide-image-container">
                        <img class="slide-image" src="{{ slide.image }}" alt="{{ slide.title }}">
                    </div>
                    {% endif %}
                    <div class="slide-content-container">
                        <h2 class="slide-title">{{ slide.title }}</h2>
                        {% if slide.content %}
                        <p class="slide-content">{{ slide.content }}</p>
                        {% endif %}
                        {% if slide.bullets %}
                        <ul class="bullets">
                            {% for bullet in slide.bullets %}
                            <li>{{ bullet }}</li>
                            {% endfor %}
                        </ul>
                        {% endif %}
                    </div>
                </section>
                {% endfor %}
            </div>
        </div>
        <nav class="navigation">
            <button class="nav-btn" id="prevBtn" aria-label="Previous Slide">←</button>
            <button class="nav-btn" id="nextBtn" aria-label="Next Slide">→</button>
        </nav>
        <div class="slide-indicator" id="slideIndicator" aria-live="polite">1 / {{ slides|length }}</div>
        <div class="progress-bar" id="progressBar" role="progressbar"></div>
    </div>

    <script>
        document.addEventListener('DOMContentLoaded', () => {
            const slides = document.getElementById('slides');
            const prevBtn = document.getElementById('prevBtn');
            const nextBtn = document.getElementById('nextBtn');
            const slideIndicator = document.getElementById('slideIndicator');
            const progressBar = document.getElementById('progressBar');
            const totalSlides = {{ slides|length }};
            let currentSlide = 0;

            const updateSlide = () => {
                slides.style.transform = `translateX(-${currentSlide * 100}%)`;
                slideIndicator.textContent = `${currentSlide + 1} / ${totalSlides}`;
                prevBtn.disabled = currentSlide === 0;
                nextBtn.disabled = currentSlide === totalSlides - 1;
                progressBar.style.width = `${((currentSlide + 1) / totalSlides) * 100}%`;
            };

            const previousSlide = () => {
                if (currentSlide > 0) {
                    currentSlide--;
                    updateSlide();
                }
            };

            const nextSlide = () => {
                if (currentSlide < totalSlides - 1) {
                    currentSlide++;
                    updateSlide();
                }
            };

            prevBtn.addEventListener('click', previousSlide);
            nextBtn.addEventListener('click', nextSlide);
            document.addEventListener('keydown', (e) => {
                if (e.key === 'ArrowLeft') previousSlide();
                if (e.key === 'ArrowRight') nextSlide();
            });
            window.addEventListener('resize', updateSlide);

            updateSlide();
        });
    </script>
</body>
</html>
"""
}

def fetch_image_from_pexels(query):
    """Fetch an image from Pexels API."""
    try:
        headers = {"Authorization": PEXELS_API_KEY}
        params = {"query": query, "per_page": 1}
        response = requests.get(PEXELS_API_URL, headers=headers, params=params)
        if response.status_code == 200:
            data = response.json()
            if data['photos']:
                img_url = data['photos'][0]['src']['large']
                img_response = requests.get(img_url)
                return io.BytesIO(img_response.content)
        logger.warning(f"No image found for query '{query}'")
        return None
    except Exception as e:
        logger.error(f"Error fetching image from Pexels: {str(e)}")
        return None

def lighten_image(img_stream, factor=0.9):
    """Lighten an image for use as a background."""
    try:
        img = Image.open(img_stream).convert("RGB")
        enhancer = Image.new("RGB", img.size, (255, 255, 255))
        img = Image.blend(img, enhancer, factor)
        new_stream = io.BytesIO()
        img.save(new_stream, format='PNG')
        new_stream.seek(0)
        return new_stream
    except Exception as e:
        logger.error(f"Error lightening image: {str(e)}")
        return img_stream

def add_shadow_to_shape(shape):
    """Add a shadow effect to a shape."""
    try:
        sp = shape._element
        spPr = sp.find('{http://schemas.openxmlformats.org/drawingml/2006/main}spPr')
        if spPr is None:
            spPr = OxmlElement('a:spPr')
            sp.append(spPr)
        
        effect_lst = OxmlElement('a:effectLst')
        outer_shdw = OxmlElement('a:outerShdw')
        outer_shdw.set('dist', '20000')
        outer_shdw.set('dir', '2700000')
        outer_shdw.set('algn', 'ctr')
        srgb_clr = OxmlElement('a:srgbClr')
        srgb_clr.set('val', '000000')
        alpha = OxmlElement('a:alpha')
        alpha.set('val', '40000')
        srgb_clr.append(alpha)
        outer_shdw.append(srgb_clr)
        effect_lst.append(outer_shdw)
        spPr.append(effect_lst)
    except Exception as e:
        logger.warning(f"Failed to apply shadow to shape: {e}")

def apply_element_properties(shape, properties):
    """Apply formatting properties to a shape."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = True
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Calibri'
            if 'font_size' in properties:
                run.font.size = Pt(properties['font_size'])
            if 'font_color' in properties:
                r, g, b = properties['font_color']
                run.font.color.rgb = PPTXRGBColor(r, g, b)
    if 'alignment' in properties:
        align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
        tf.paragraphs[0].alignment = align_map.get(properties['alignment'], PP_ALIGN.LEFT)
    if properties.get('shadow', False):
        try:
            add_shadow_to_shape(shape)
        except Exception as e:
            logger.warning(f"Failed to apply shadow: {e}")

def add_custom_image(slide, img_stream, properties):
    """Add an image to a slide with specified properties."""
    try:
        left = Inches(properties.get('position', [0, 0])[0])
        top = Inches(properties.get('position', [0, 0])[1])
        width = Inches(properties.get('size', [6, 4])[0])
        height = Inches(properties.get('size', [6, 4])[1])
        slide.shapes.add_picture(img_stream, left, top, width=width, height=height)
    except Exception as e:
        logger.error(f"Error adding custom image: {str(e)}")

def add_slide(prs, slide_data, slide_config, topic, background_img_data=None):
    """Add a formatted slide to the presentation with dynamic positioning."""
    try:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        if background_img_data:
            light_stream = lighten_image(io.BytesIO(background_img_data))
            add_custom_image(slide, light_stream, {
                "position": [0, 0],
                "size": [10, 7.5]
            })

        header_shape = slide.shapes.add_shape(
            1,
            Inches(0), Inches(0),
            Inches(10), Inches(slide_config['header']['height'])
        )
        fill = header_shape.fill
        fill.solid()
        r, g, b = slide_config['header']['color']
        fill.fore_color.rgb = PPTXRGBColor(r, g, b)

        title_text = slide_data.get('title', 'Untitled')
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.0))
        tf = title_box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title_text
        apply_element_properties(title_box, slide_config['title'])

        content_y = 1.8
        max_bullets = min(len(slide_data.get('bullets', [])), 3)

        for i, item in enumerate(slide_data.get('bullets', [])[:max_bullets]):
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(content_y), Inches(9.0), Inches(0.8)
            )
            tf = content_box.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.auto_size = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = f"• {item}"
            apply_element_properties(content_box, slide_config['content'])
            content_y += 1.0 + (len(item) // 80) * 0.4

        content = slide_data.get('bullets', [])
        icon_keyword = (content[0].split()[0].lower() if content else topic.lower().split()[0])
        icon_stream = fetch_image_from_pexels(icon_keyword)
        if icon_stream:
            footer_y = max(6.0, content_y + 0.5)
            footer_config = slide_config['footer'].copy()
            footer_config['position'] = [6.5, footer_y]
            add_custom_image(slide, icon_stream, footer_config)

        return slide
    except Exception as e:
        logger.error(f"Error adding slide: {str(e)}")
        raise

def preprocess_yaml_content(yaml_content):
    """Preprocess YAML content to fix bullet characters and indentation."""
    try:
        if not isinstance(yaml_content, str):
            logger.error(f"Expected string for yaml_content, got {type(yaml_content)}")
            return None

        # Remove any extra document markers
        yaml_content = re.sub(r'^---\s*$', '', yaml_content, flags=re.MULTILINE)
        yaml_content = yaml_content.strip()
        
        # Fix bullet points
        yaml_content = re.sub(r'^\s*[\•*]\s+', '  - ', yaml_content, flags=re.MULTILINE)
        
        lines = yaml_content.splitlines()
        cleaned_lines = []
        for line in lines:
            stripped = line.rstrip()
            if stripped:
                leading_spaces = len(line) - len(line.lstrip())
                if stripped.lstrip().startswith('-'):
                    cleaned_lines.append(' ' * (leading_spaces - (leading_spaces % 2)) + stripped.lstrip())
                else:
                    cleaned_lines.append(line)
            else:
                cleaned_lines.append(line)
        
        cleaned_yaml = '\n'.join(cleaned_lines)
        logger.debug(f"Preprocessed YAML content:\n{cleaned_yaml}")
        return cleaned_yaml
    except Exception as e:
        logger.error(f"Error preprocessing YAML content: {str(e)}")
        return None

def create_pptx_from_yaml(yaml_content, output_path, topic, email):
    """Create PowerPoint presentation from YAML content and save both to disk and DB."""
    try:
        if not isinstance(yaml_content, str):
            logger.error(f"Expected string for yaml_content, got {type(yaml_content)}")
            return {"success": False, "error": f"Expected string for yaml_content, got {type(yaml_content)}"}
        if not yaml_content.strip():
            logger.error("yaml_content is empty")
            return {"success": False, "error": "yaml_content is empty"}
        if not output_path.endswith('.pptx'):
            logger.error("Output path must end with .pptx")
            return {"success": False, "error": "Output path must end with .pptx"}

        os.makedirs(os.path.dirname(output_path), exist_ok=True)

        cleaned_yaml = preprocess_yaml_content(yaml_content)
        if not cleaned_yaml:
            logger.error("Failed to preprocess YAML content")
            return {"success": False, "error": "Failed to preprocess YAML content"}

        try:
            data = yaml.safe_load(cleaned_yaml)
            if data is None:
                documents = list(yaml.safe_load_all(cleaned_yaml))
                if documents:
                    data = documents[0]
                else:
                    raise yaml.YAMLError("Empty YAML content")
        except yaml.YAMLError as e:
            logger.error(f"YAML parsing error: {str(e)}")
            return {"success": False, "error": f"YAML parsing error: {str(e)}"}

        if not data or 'presentation' not in data:
            logger.error("Invalid YAML structure: missing 'presentation' key")
            return {"success": False, "error": "Invalid YAML structure: missing 'presentation' key"}
        
        slides_data = data.get('presentation', {}).get('slides', [])
        if not slides_data:
            logger.error("No slides found in YAML")
            return {"success": False, "error": "No slides found"}

        include_images = data.get('presentation', {}).get('include_images', True)

        prs = Presentation()

        # Your slide creation code here (title slide and others) ...
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_shape = title_slide.shapes.title
        title_shape.text = data.get('presentation', {}).get('title', topic)
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.color.rgb = PPTXRGBColor(0, 51, 102)
        subtitle = title_slide.placeholders[1]
        subtitle.text = f"Exploring {topic}"
        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle.text_frame.paragraphs[0].font.color.rgb = PPTXRGBColor(50, 50, 50)

        background_img_data = fetch_consistent_background_image(topic)

        slide_config = {
            "header": {
                "color": [0, 51, 102],
                "height": 1.2
            },
            "title": {
                "font_size": 32,
                "font_color": [255, 255, 255],
                "alignment": "center",
                "shadow": True
            },
            "content": {
                "font_size": 22,
                "font_color": [10, 10, 10],
                "position": [0.5, 2.0],
                "shadow": True
            },
            "footer": {
                "color": [100, 100, 100],
                "position": [6.5, 4.0],
                "size": [3.0, 3.0]
            }
        }

        for slide_data in slides_data:
            if include_images:
                img_url = search_pexels_image(topic, slide_title=slide_data.get('title'))
                if img_url:
                    img_data = download_image(img_url)
                    if img_data:
                        img_stream = io.BytesIO(img_data)
                        slide = prs.slides.add_slide(prs.slide_layouts[5])
                        slide.shapes.title.text = f"{slide_data.get('title', 'Untitled')} - Visual"
                        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = PPTXRGBColor(245, 245, 220)
                        slide.shapes.title.text_frame.paragraphs[0].font.name = 'Arial'
                        img = Image.open(img_stream)
                        width, height = img.size
                        max_w, max_h = Inches(10), Inches(5.5)
                        scale = min(max_w / width, max_h / height)
                        new_w, new_h = width * scale, height * scale
                        left = (prs.slide_width - new_w) / 2
                        top = (prs.slide_height - new_h + Inches(1)) / 2
                        slide.shapes.add_picture(io.BytesIO(img_data), left, top, width=new_w, height=new_h)
            add_slide(prs, slide_data, slide_config, topic, background_img_data)

        # Save PPTX file to disk (existing functionality)
        prs.save(output_path)
        logger.info(f"PowerPoint saved to {output_path}")

        # Also save PPTX as base64 string in DB
        pptx_stream = io.BytesIO()
        prs.save(pptx_stream)
        pptx_stream.seek(0)
        base64_pptx = base64.b64encode(pptx_stream.read()).decode('utf-8')

        file_record = FileRecord(
            topic=topic,
            user_email=email,
            file_type='pptx',
            file_data=base64_pptx
        )
        db.session.add(file_record)
        db.session.commit()
        logger.info(f"PowerPoint stored in DB for email: {email}, topic: {topic}")

        return {"success": True, "file_id": file_record.id}

    except Exception as e:
        logger.error(f"Error in create_pptx_from_yaml: {str(e)}")
        return {"success": False, "error": str(e)}

def create_docx_from_yaml(yaml_content, output_path, email):
    """Create Word document from YAML content and store as base64 in DB."""
    try:
        if not isinstance(yaml_content, str):
            return {"success": False, "error": f"Expected string for yaml_content, got {type(yaml_content)}"}
        if not yaml_content.strip():
            return {"success": False, "error": "yaml_content is empty"}
        if not output_path.endswith('.docx'):
            return {"success": False, "error": "Output path must end with .docx"}

        os.makedirs(os.path.dirname(output_path), exist_ok=True)

        cleaned_yaml = preprocess_yaml_content(yaml_content)
        if not cleaned_yaml:
            logger.error("Failed to preprocess YAML content")
            return {"success": False, "error": "Failed to preprocess YAML content"}

        try:
            data = yaml.safe_load(cleaned_yaml)
            if data is None:
                documents = list(yaml.safe_load_all(cleaned_yaml))
                if documents:
                    data = documents[0]
                else:
                    raise yaml.YAMLError("Empty YAML content")
        except yaml.YAMLError as e:
            logger.error(f"YAML parsing error: {str(e)}")
            return {"success": False, "error": f"YAML parsing error: {str(e)}"}

        slides = data.get('presentation', {}).get('slides', [])
        if not slides:
            return {"success": False, "error": "No slides found"}

        # Generate Word document
        doc = Document()
        doc.add_heading('Presentation', 0)
        for slide in slides:
            doc.add_heading(slide.get('title', 'Untitled'), level=1)
            for bullet in slide.get('bullets', []):
                doc.add_paragraph(bullet, style='ListBullet')

        # Save the file locally
        doc.save(output_path)
        logger.info(f"Word document saved to {output_path}")

        # Read file and encode in base64
        with open(output_path, "rb") as f:
            encoded_data = base64.b64encode(f.read()).decode('utf-8')

        # Extract topic from YAML (fallback to filename if not found)
        topic = data.get('presentation', {}).get('title') or os.path.basename(output_path).split('.')[0]

        # Store record in DB
        file_record = FileRecord(
            topic=topic,
            user_email=email,
            file_type='docx',
            file_data=encoded_data
        )
        db.session.add(file_record)
        db.session.commit()
        logger.info(f"FileRecord created in DB for {email}")

        return {"success": True, "file_record_id": file_record.id}
    
    except Exception as e:
        logger.error(f"Error in create_docx_from_yaml: {str(e)}")
        return {"success": False, "error": str(e)}

def create_pdf_from_yaml(yaml_content, output_path, email):
    """Create PDF document from YAML content and store as base64 in DB."""
    try:
        if not isinstance(yaml_content, str):
            return {"success": False, "error": f"Expected string for yaml_content, got {type(yaml_content)}"}
        if not yaml_content.strip():
            return {"success": False, "error": "yaml_content is empty"}
        if not output_path.endswith('.pdf'):
            return {"success": False, "error": "Output path must end with .pdf"}

        os.makedirs(os.path.dirname(output_path), exist_ok=True)

        cleaned_yaml = preprocess_yaml_content(yaml_content)
        if not cleaned_yaml:
            logger.error("Failed to preprocess YAML content")
            return {"success": False, "error": "Failed to preprocess YAML content"}

        try:
            data = yaml.safe_load(cleaned_yaml)
            if data is None:
                documents = list(yaml.safe_load_all(cleaned_yaml))
                if documents:
                    data = documents[0]
                else:
                    raise yaml.YAMLError("Empty YAML content")
        except yaml.YAMLError as e:
            logger.error(f"YAML parsing error: {str(e)}")
            return {"success": False, "error": f"YAML parsing error: {str(e)}"}

        slides = data.get('presentation', {}).get('slides', [])
        if not slides:
            return {"success": False, "error": "No slides found"}

        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=14)
        pdf.cell(0, 10, txt="Presentation", ln=True, align='C')
        pdf.ln(10)

        for slide in slides:
            pdf.set_font("Arial", 'B', 12)
            pdf.cell(0, 10, txt=slide.get('title', 'Untitled Slide'), ln=True)
            pdf.set_font("Arial", size=11)
            for bullet in slide.get('bullets', []):
                pdf.set_x(20)
                pdf.multi_cell(0, 8, txt=f"- {bullet}")
                pdf.ln(2)
            pdf.ln(5)

        # Save locally
        pdf.output(output_path)
        logger.info(f"PDF saved to {output_path}")

        # Read file as base64
        with open(output_path, "rb") as f:
            encoded_data = base64.b64encode(f.read()).decode('utf-8')

        topic = data.get('presentation', {}).get('title') or os.path.basename(output_path).split('.')[0]

        # Store in DB
        file_record = FileRecord(
            topic=topic,
            user_email=email,
            file_type='pdf',
            file_data=encoded_data
        )
        db.session.add(file_record)
        db.session.commit()
        logger.info(f"PDF FileRecord created for {email}")

        return {"success": True, "file_record_id": file_record.id}

    except Exception as e:
        logger.error(f"Error in create_pdf_from_yaml: {str(e)}")
        return {"success": False, "error": str(e)}

def process_images_for_slides(yaml_data, topic):
    """Process images for slides, adding base64-encoded data URIs."""
    try:
        if not isinstance(yaml_data, dict) or 'presentation' not in yaml_data:
            logger.error("Invalid YAML data: missing 'presentation' key")
            return yaml_data

        include_images = yaml_data.get('presentation', {}).get('include_images', True)
        if not include_images:
            logger.info("Images disabled for presentation")
            return yaml_data

        slides = yaml_data['presentation'].get('slides', [])
        if not slides:
            logger.warning("No slides found in YAML data")
            return yaml_data

        for slide in slides:
            if slide.get('needs_image', True):
                # Use slide title as primary keyword, fall back to topic
                keyword = slide.get('title', topic)
                img_url = search_pexels_image(topic, slide_title=keyword)
                if img_url:
                    img_data = download_image(img_url)
                    if img_data:
                        # Determine image format using PIL
                        img = Image.open(io.BytesIO(img_data))
                        format_map = {'JPEG': 'image/jpeg', 'PNG': 'image/png'}
                        mime_type = format_map.get(img.format, 'image/jpeg')
                        # Convert image data to base64 data URI
                        img_base64 = base64.b64encode(img_data).decode('utf-8')
                        slide['image'] = f"data:{mime_type};base64,{img_base64}"
                        logger.info(f"Added image for slide '{slide.get('title', 'Untitled')}'")
                    else:
                        logger.warning(f"Failed to download image for keyword: {keyword}")
                else:
                    logger.warning(f"No image found for keyword: {keyword}")
        return yaml_data
    except Exception as e:
        logger.error(f"Error processing images for slides: {str(e)}")
        return yaml_data

def create_html_from_yaml(yaml_content, output_path, topic, email, html_presentation_type='minimalist'):
    """Create HTML presentation from YAML content and store base64 in DB."""
    try:
        if not isinstance(yaml_content, str):
            logger.error(f"Expected string for yaml_content, got {type(yaml_content)}")
            return {"success": False, "error": f"Expected string for yaml_content, got {type(yaml_content)}"}
        if not yaml_content.strip():
            logger.error("yaml_content is empty")
            return {"success": False, "error": "yaml_content is empty"}
        if not output_path.endswith('.html'):
            logger.error("Output path must end with .html")
            return {"success": False, "error": "Output path must end with .html"}

        os.makedirs(os.path.dirname(output_path), exist_ok=True)

        cleaned_yaml = preprocess_yaml_content(yaml_content)
        if not cleaned_yaml:
            logger.error("Failed to preprocess YAML content")
            return {"success": False, "error": "Failed to preprocess YAML content"}

        try:
            data = yaml.safe_load(cleaned_yaml)
            if data is None:
                documents = list(yaml.safe_load_all(cleaned_yaml))
                if documents:
                    data = documents[0]
                else:
                    raise yaml.YAMLError("Empty YAML content")
        except yaml.YAMLError as e:
            logger.error(f"YAML parsing error: {str(e)}")
            return {"success": False, "error": f"YAML parsing error: {str(e)}"}

        if not data or 'presentation' not in data:
            logger.error("Invalid YAML structure: missing 'presentation' key")
            return {"success": False, "error": "Invalid YAML structure: missing 'presentation' key"}

        presentation_data = data.get('presentation', {})
        slides = presentation_data.get('slides', [])
        if not slides:
            logger.error("No slides found in YAML")
            return {"success": False, "error": "No slides found"}

        # Process images
        data = process_images_for_slides(data, topic)

        # Prepare data for Jinja template
        template_data = {
            'title': presentation_data.get('title', 'AI-Generated Presentation'),
            'slides': [
                {
                    'title': slide.get('title', 'Untitled Slide'),
                    'content': slide.get('content', ''),
                    'bullets': slide.get('bullets', []),
                    'image': slide.get('image', '')
                } for slide in slides if slide.get('title') or slide.get('content') or slide.get('bullets')
            ]
        }

        # Render HTML
        jinja_template = Template(HTML_TEMPLATES.get(html_presentation_type, HTML_TEMPLATES['minimalist']))
        html_content = jinja_template.render(**template_data)

        # Save to file
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(html_content)

        logger.info(f"HTML presentation saved to {output_path}")

        # Encode content as base64
        encoded_data = base64.b64encode(html_content.encode('utf-8')).decode('utf-8')

        # Store in DB
        file_record = FileRecord(
            topic=topic,
            user_email=email,
            file_type='html',
            file_data=encoded_data
        )
        db.session.add(file_record)
        db.session.commit()
        logger.info(f"HTML FileRecord created for {email}")

        return {"success": True, "file_record_id": file_record.id}

    except Exception as e:
        logger.error(f"Error in create_html_from_yaml: {str(e)}")
        return {"success": False, "error": str(e)}